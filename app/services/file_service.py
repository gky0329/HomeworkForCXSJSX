from pathlib import Path

from app.services.i18n import tr


def _stripped(text: str) -> str:
    return "\n".join(
        line for line in text.splitlines()
        if line.strip()
    )


def _extract_pdf(path: Path) -> str:
    import fitz
    doc = fitz.open(str(path))
    try:
        lines: list[str] = []
        for page in doc:
            blocks = page.get_text("blocks")
            blocks.sort(key=lambda b: (b[1], b[0]))
            for b in blocks:
                text = b[4].strip()
                if not text:
                    continue
                y_ratio = b[1] / page.rect.height
                if y_ratio < 0.05 or y_ratio > 0.93:
                    continue
                lines.append(text)
        return "\n\n".join(lines)
    finally:
        doc.close()


def _extract_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paras)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides_out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        if len(lines) > 1:
            slides_out.append("\n".join(lines))
    return "\n\n---\n\n".join(slides_out)


def _extract_markdown(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    lines = []
    in_code = False
    for line in raw.splitlines():
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            lines.append(line)
            continue
        if stripped.startswith("#"):
            lines.append(stripped.lstrip("# "))
            continue
        if stripped.startswith("- ") or stripped.startswith("* "):
            lines.append(stripped[2:])
            continue
        if stripped.startswith("> "):
            lines.append(stripped[2:])
            continue
        if stripped:
            lines.append(stripped)
    return "\n".join(lines)


def _extract_code(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    return raw


def _extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


HANDLERS: dict[str, list[str]] = {
    "pdf":  [".pdf"],
    "docx": [".docx", ".doc"],
    "pptx": [".pptx", ".ppt"],
    "markdown": [".md", ".mdx"],
    "code":  [".cpp", ".cxx", ".cc", ".c", ".h", ".hpp", ".py", ".java", ".js", ".ts"],
    "text":  [".txt"],
}


def _handler_for(ext: str) -> str | None:
    for name, exts in HANDLERS.items():
        if ext in exts:
            return name
    return None


_EXTRACTORS = {
    "pdf":      _extract_pdf,
    "docx":     _extract_docx,
    "pptx":     _extract_pptx,
    "markdown": _extract_markdown,
    "code":     _extract_code,
    "text":     _extract_text,
}


SUPPORTED_EXTENSIONS = sorted(
    ext for exts in HANDLERS.values() for ext in exts
)


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    ext = path.suffix.lower()
    handler = _handler_for(ext)
    if handler is None:
        raise ValueError(
            tr(
                "Unsupported file type: {ext}. Supported: {supported}",
                ext=ext,
                supported=", ".join(SUPPORTED_EXTENSIONS),
            )
        )
    return _EXTRACTORS[handler](path)


def file_type_label(ext: str) -> str:
    labels = {
        ".pdf": tr("PDF Document"),
        ".docx": tr("Word Document"),
        ".doc": tr("Word Document"),
        ".md": "Markdown",
        ".mdx": "Markdown",
        ".cpp": tr("C++ Source"),
        ".cxx": tr("C++ Source"),
        ".cc": tr("C++ Source"),
        ".c": tr("C Source"),
        ".h": tr("C/C++ Header"),
        ".hpp": tr("C++ Header"),
        ".py": tr("Python Source"),
        ".java": tr("Java Source"),
        ".js": tr("JavaScript"),
        ".ts": tr("TypeScript"),
        ".txt": tr("Plain Text"),
        ".pptx": tr("PowerPoint"),
        ".ppt": tr("PowerPoint"),
    }
    return labels.get(ext, tr("{ext} File", ext=ext))
