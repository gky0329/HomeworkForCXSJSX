"""
生成演示 PPTX（中文，12页，7分钟，匹配用户讲稿）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色
BG     = RGBColor(0x1E, 0x1E, 0x1E)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
WHITE  = RGBColor(0xD4, 0xD4, 0xD4)
GRAY   = RGBColor(0x80, 0x80, 0x80)
STACK  = RGBColor(0x56, 0x9C, 0xD6)
HEAP   = RGBColor(0xCE, 0x91, 0x78)
GREEN  = RGBColor(0x4E, 0xC9, 0xB0)
RED    = RGBColor(0xF4, 0x47, 0x47)
DARK2  = RGBColor(0x2D, 0x2D, 0x2D)
BORDER = RGBColor(0x3E, 0x3E, 0x3E)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

W = Inches(13.333)
H = Inches(7.5)
OUT = "docs/C++_Memory_Visualizer.pptx"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def S(slide, c=BG):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = c

def T(slide, l, t, w, h, text, color=WHITE, size=18, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def H1(slide, text):
    T(slide, 1.0, 0.4, 11.3, 0.7, text, WHITE, 32, True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.1), Inches(0.7), Pt(4))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()

def B(slide, items, y=1.6, size=16, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.0), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = "Arial"; p.space_after = Pt(10)

def VID(slide, label):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(1.8), Inches(9.3), Inches(3.8))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = BORDER; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "▶  " + label; p.font.size = Pt(22)
    p.font.color.rgb = ACCENT; p.font.bold = True; p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "【此处插入演示视频】"; p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY; p2.font.name = "Arial"; p2.alignment = PP_ALIGN.CENTER

def FTR(slide, n):
    T(slide, 1.0, 6.9, 11.3, 0.3, f"{n} / 12", GRAY, 10, align=PP_ALIGN.RIGHT)

def CARD(slide, l, t, tt, desc, color=STACK):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(5.4), Inches(1.7))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = BORDER; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.text = tt; p.font.size = Pt(15); p.font.color.rgb = color; p.font.bold = True; p.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(13); p2.font.color.rgb = WHITE; p2.font.name = "Arial"; p2.space_before = Pt(4)

def CODE(slide, y, text):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(y), Inches(9.3), Inches(0.7))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x25,0x25,0x26); s.line.color.rgb = BORDER; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = "Courier New"

N = {}

# ═══════════════════════════════════════════
# SLIDE 1 封面
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s)
T(s, 1.0, 2.0, 11.3, 1.0, "C++ 工作台", WHITE, 48, True, PP_ALIGN.CENTER)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.9), Inches(2.3), Pt(4))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
T(s, 1.0, 3.1, 11.3, 0.6, "实用至上 · 用户体验至上", GRAY, 22, align=PP_ALIGN.CENTER)
T(s, 1.0, 4.2, 11.3, 0.5, "全流程 · 可视化 · 可持续 · 可循环", WHITE, 16, align=PP_ALIGN.CENTER)
T(s, 1.0, 5.4, 11.3, 0.4, "AI 辅助 C++ 学习工具  |  路演发布会", GRAY, 14, align=PP_ALIGN.CENTER)
N[0] = "大家好，欢迎来到我们的大作业项目C++工作台的路演发布会。首先，我不躲不藏不绕，用一句话概括核心理念：实用至上，用户体验至上。"

# ═══════════════════════════════════════════
# SLIDE 2 核心理念
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "AI 不是敌人，是工具")
B(s, [
    "AI 有不足：提示词繁琐 · 长期记忆差 · 局限于文本模态 · 交互体验不理想",
    "",
    "我们的思路：",
    "搭建一个 全流程 、可视化 、可持续 、可循环 的高效 AI 辅助学习工具",
    "实现超越传统 Agent 的交互式学习体验",
], y=1.6, size=18)
T(s, 1.2, 4.8, 11.0, 0.8, "「站在 AI 的对立面是没有意义的。\n利用 AI，才能如虎添翼。」", GRAY, 16, False)
FTR(s, 2)
N[1] = "AI时代，人工智能本身就可以辅助学习。但AI有很多不足：写提示词繁琐，长期记忆差，局限于文本模态。因此我们工作的核心思路是搭建全流程、可视化、可持续、可循环的AI辅助学习工具，实现超越传统agent的交互式学习体验。"

# ═══════════════════════════════════════════
# SLIDE 3 主界面 + 设置
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "主界面 & 设置")
B(s, [
    "🏠  主界面：功能简洁明了，统计信息帮助掌握学习进度",
    "",
    "⚙  设置：自由选择 AI 网址、模型、API Key",
    "       支持 DeepSeek / OpenAI / Claude / Gemini 各大厂商",
    "       首次设置后长期有效，无需每次填写",
    "",
    "🎨  前端借鉴知名游戏设计风格，后续提供极简风格选项",
], y=1.6, size=18)
FTR(s, 3)
N[2] = "运行程序直接进入主界面，功能简洁明了。统计信息便于掌握学习进度。设置中可以选择AI网址、模型和API Key，首次设置长期有效。前端借鉴某知名游戏设计风格。"

# ═══════════════════════════════════════════
# SLIDE 4 代码编辑器（核心）
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "可视化代码编辑器 — 核心设计")
B(s, [
    "📝  实时呈现内存状态：变量 · 指针 · 内存 · 对象的关系一目了然",
    "",
    "🚀  内置示例代码，一键运行",
    "       点击自动播放，查看内存逐行变化过程",
    "",
    "🎨  精妙画布设计：缩放 · 移动丝滑流畅",
    "       即使多层嵌套也没有问题",
], y=1.6, size=18)
CODE(s, 4.5, "int a = 42;   int* p = new int(100);   int* q = &a;   *p = 200;   delete p;")
FTR(s, 4)
N[3] = "最核心的设计：可视化代码编辑器。实时呈现内存状态。这里给出示例代码，点击运行生成可视化图形，自动播放查看内存变化。画布精妙设计，缩放移动丝滑流畅。"

# ═══════════════════════════════════════════
# SLIDE 5 🎬 演示① 代码编辑器
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "🎬 演示 — 代码编辑器操作")
VID(s, "写代码 → Run → Auto Play → Canvas 动画 (约 60s)")
B(s, ["展示：运行代码 · 自动播放 · 内存变化 · 画布交互"], y=5.8, size=14, color=GRAY)
FTR(s, 5)
N[4] = "【插入视频】展示代码编辑器完整操作流程。"

# ═══════════════════════════════════════════
# SLIDE 6 OJ 分析
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "OJ 分析")
B(s, [
    "💡  解决 OJ 题目时，如何获得详细指导？",
    "",
    "✅  点击分析 → 生成示例代码 → 可调用本地编译器测试",
    "✅  相关思路讲解 · 知识点标注 · 一键加入复习",
    "✅  算法复杂度分析 · 常见错误提示",
    "✅  一键发送代码到可视化区域",
], y=1.6, size=18)
FTR(s, 6)
N[5] = "OJ分析部分。粘贴题目后点击分析，生成示例代码，可本地编译测试。还有思路讲解、知识点标注、一键加入复习，非常全面方便。"

# ═══════════════════════════════════════════
# SLIDE 7 文件导入
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "文件导入")
B(s, [
    "📂  快速分析代码/讲义/课件，支持多种格式（PDF · DOCX · PPTX · MD · CPP）",
    "",
    "🤖  上传文件 → 提取知识点 → AI 拆分总结",
    "       知识点可 可视化 或 加入知识库",
    "",
    "📝  配套基础练习题，实时检验学习成果",
    "       错题一键加入错题集",
], y=1.6, size=18)
FTR(s, 7)
N[6] = "文件导入部分提供快速分析功能。支持多种格式。上传PDF后点击提取知识点，AI总结拆分。配套练习题实时检验，错题加入错题集。"

# ═══════════════════════════════════════════
# SLIDE 8 🎬 演示② OJ + FileImport
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "🎬 演示 — OJ 分析 & 文件导入")
VID(s, "OJ 分析 · 文件导入 · 知识点提取 · 在线测试 (约 60s)")
B(s, ["展示：OJ 分析 → 知识点标注 → 文件导入 → AI 提取 → 在线测试"], y=5.8, size=14, color=GRAY)
FTR(s, 8)
N[7] = "【插入视频】展示 OJ 分析和文件导入流程。"

# ═══════════════════════════════════════════
# SLIDE 9 错题复习
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "错题复习 — 可持续 · 可循环的核心")
B(s, [
    "🃏  错题自动分类储存，可按分类加载复习",
    "",
    "💡  AI 提示：不确定时，AI 在不给出答案的前提下引导思考",
    "📖  显示答案：正确答案 + AI 详细讲解",
    "",
    "📝  记录学习心得 · 评价题目难度",
    "       难度评价决定以后看到这道题的频率（SM-2 算法）",
    "",
    "➕  也可手动添加错题，统一整理复习",
], y=1.6, size=16)
FTR(s, 9)
N[8] = "错题复习是可持续可循环的核心支撑。错题自动分类。每道题有提示和显示答案。提示由AI在不给出答案的前提下引导。显示答案后可以记录心得、评价难度，这决定了后续复习频率。也可以手动添加错题。"

# ═══════════════════════════════════════════
# SLIDE 10 知识库
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "知识库 — 总体的储存库")
B(s, [
    "📚  列表界面：清晰列出所有入库知识点",
    "",
    "✨  点开知识点：AI 详细讲解（Markdown 渲染）",
    "       一键加入复习 · 一键生成小测验",
    "       按需删除整理",
    "",
    "🗺  图谱界面：知识点聚集显示，可点击查看",
    "       力导向图，错误越多节点越大",
], y=1.6, size=17)
FTR(s, 10)
N[9] = "知识库是总体的储存库。列表界面清晰列出知识点。点开后看到AI讲解，可以一键加入复习、一键生成小测验。图谱界面聚集显示知识点，点击查看。"

# ═══════════════════════════════════════════
# SLIDE 11 🎬 演示③ Review + KB
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s); H1(s, "🎬 演示 — 错题复习 & 知识库")
VID(s, "错题复习 · 分类加载 · AI 提示 · 知识库 · 小测验 (约 60s)")
B(s, ["展示：错题分类 · AI 提示 · 显示答案 · 难度评价 · 知识库讲解 · 小测验"], y=5.8, size=14, color=GRAY)
FTR(s, 11)
N[10] = "【插入视频】展示错题复习和知识库的完整交互流程。"

# ═══════════════════════════════════════════
# SLIDE 12 总结
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); S(s)
T(s, 1.0, 2.0, 11.3, 1.0, "不做花里胡哨", WHITE, 44, True, PP_ALIGN.CENTER)
T(s, 1.0, 2.8, 11.3, 0.8, "做一款真的懂你的小程序", WHITE, 28, False, PP_ALIGN.CENTER)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(2.3), Pt(4))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
B(s, [
    "实用至上 · 用户体验至上",
    "全流程 · 可视化 · 可持续 · 可循环",
    "如果能在实际学习中有所帮助，便是最大的成就",
], y=4.0, size=18)
T(s, 1.0, 5.6, 11.3, 0.4, "github.com/gky0329/HomeworkForCXSJSX", ACCENT, 14, align=PP_ALIGN.CENTER)
T(s, 1.0, 6.2, 11.3, 0.3, "谢谢大家", GRAY, 16, align=PP_ALIGN.CENTER)
FTR(s, 12)
N[11] = "总得来看，我们的项目并不追求花里胡哨的功能，而是从实际体验出发，做出一款真的懂你的小程序。如果能在实际有所帮助，便是最大的成就。谢谢大家。"

# ═══ 演讲备注 ═══
for i, note in N.items():
    if i < len(prs.slides):
        prs.slides[i].notes_slide.notes_text_frame.text = note

prs.save(OUT)
print(f"✅ {OUT}")
print(f"   共 12 页  |  3 个视频位  |  约 7 分钟")
